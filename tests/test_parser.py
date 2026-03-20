import os
import shutil
import sys
import tempfile
import unittest
import warnings
from pathlib import Path

from docx import Document
from pptx import Presentation
from pypdf import PdfWriter

# 添加 deepdoc 路径到 sys.path
sys.path.append(os.path.join(os.path.dirname(__file__), ".."))

# 导入所有解析器（使用正确的类名）
from deepdoc import TxtParser, MarkdownParser, JsonParser, HtmlParser, ExcelParser, PptParser, DocxParser, PdfParser


class TestParsers(unittest.TestCase):
    """测试所有解析器的功能"""

    def setUp(self):
        """测试前的准备工作"""
        # 创建测试数据目录
        self.test_data_dir = Path(__file__).parent / "test_data"
        self.test_data_dir.mkdir(exist_ok=True)

        # 创建测试结果目录
        self.test_results_dir = Path(__file__).parent / "test_results"
        self.test_results_dir.mkdir(exist_ok=True)

        # 创建临时文件
        self.temp_dir = tempfile.mkdtemp()

    def tearDown(self):
        """测试后的清理工作"""
        # 清理临时文件
        shutil.rmtree(self.temp_dir, ignore_errors=True)

    def create_test_file(self, content, filename, encoding="utf-8"):
        """创建测试文件"""
        file_path = self.test_data_dir / filename
        with open(file_path, "w", encoding=encoding) as f:
            f.write(content)
        return os.fspath(file_path)

    def create_test_docx_file(self, filename):
        """创建测试Word文件"""
        file_path = Path(self.temp_dir) / filename
        doc = Document()
        doc.add_heading("测试Word标题", level=1)
        doc.add_paragraph("这是Word正文内容。")
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "列1"
        table.cell(0, 1).text = "列2"
        table.cell(1, 0).text = "数据1"
        table.cell(1, 1).text = "数据2"
        doc.save(file_path)
        return os.fspath(file_path)

    def create_test_pptx_file(self, filename):
        """创建测试PPT文件"""
        file_path = Path(self.temp_dir) / filename
        presentation = Presentation()

        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "测试PPT标题"
        slide.placeholders[1].text = "第一页内容"

        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "第二页标题"
        slide.placeholders[1].text = "第二页内容"

        presentation.save(file_path)
        return os.fspath(file_path)

    def create_test_pdf_file(self, filename):
        """创建测试PDF文件"""
        file_path = Path(self.temp_dir) / filename
        writer = PdfWriter()
        writer.add_blank_page(width=612, height=792)
        with open(file_path, "wb") as f:
            writer.write(f)
        return os.fspath(file_path)

    def test_txt_parser(self):
        """测试文本解析器"""
        # 创建测试文本文件
        test_content = """这是一个测试文档。
包含多行文本内容。
测试中文和English混合内容。
包含特殊字符：!@#$%^&*()
包含数字：1234567890"""

        test_file = self.create_test_file(test_content, "test.txt")

        # 测试解析
        parser = TxtParser()
        result = parser(test_file)

        # 验证结果
        self.assertIsNotNone(result)
        self.assertIsInstance(result, list)
        self.assertGreater(len(result), 0)

        # 检查内容
        all_text = " ".join([chunk[0] for chunk in result])
        self.assertIn("测试文档", all_text)
        self.assertIn("中文和English", all_text)

        print(f"✅ TXT解析成功: {len(result)} 个chunk")

    def test_markdown_parser(self):
        """测试Markdown解析器"""
        test_content = """# 测试标题

这是一个**粗体**文本和*斜体*文本。

## 二级标题

- 列表项1
- 列表项2
- 列表项3

### 代码块
```python
print("Hello World")
```

> 这是一个引用块

| 表格 | 列1 | 列2 |
|------|-----|-----|
| 行1  | 数据1 | 数据2 |
| 行2  | 数据3 | 数据4 |
"""

        test_file = self.create_test_file(test_content, "test.md")

        # 测试解析
        parser = MarkdownParser()

        # 读取文件内容
        with open(test_file, "r", encoding="utf-8") as f:
            content = f.read()

        # 测试表格提取
        text, tables = parser.extract_tables_and_remainder(content)

        # 验证结果
        self.assertIsNotNone(text)
        self.assertIsInstance(tables, list)
        self.assertIn("测试标题", text)
        self.assertIn("粗体", text)

        print(f"✅ Markdown解析成功: {len(tables)} 个表格")

    def test_json_parser(self):
        """测试JSON解析器"""
        test_content = """{
    "name": "测试用户",
    "age": 25,
    "email": "test@example.com",
    "hobbies": ["读书", "游泳", "编程"],
    "address": {
        "city": "北京",
        "street": "测试街道",
        "postcode": "100000"
    },
    "active": true,
    "score": 95.5
}"""

        test_file = self.create_test_file(test_content, "test.json")

        # 读取文件为二进制
        with open(test_file, "rb") as f:
            binary_content = f.read()

        # 测试解析
        parser = JsonParser()
        result = parser(binary_content)

        # 验证结果
        self.assertIsNotNone(result)
        self.assertIsInstance(result, list)
        self.assertGreater(len(result), 0)

        # 检查内容
        all_text = " ".join(result)
        self.assertIn("测试用户", all_text)
        self.assertIn("北京", all_text)

        print(f"✅ JSON解析成功: {len(result)} 个chunk")

    def test_html_parser(self):
        """测试HTML解析器"""
        test_content = """<!DOCTYPE html>
<html>
<head>
    <title>测试页面</title>
    <meta charset="utf-8">
</head>
<body>
    <h1>主标题</h1>
    <p>这是一个<strong>段落</strong>，包含<em>强调</em>文本。</p>
    
    <h2>二级标题</h2>
    <ul>
        <li>列表项1</li>
        <li>列表项2</li>
        <li>列表项3</li>
    </ul>
    
    <table border="1">
        <tr>
            <th>表头1</th>
            <th>表头2</th>
        </tr>
        <tr>
            <td>数据1</td>
            <td>数据2</td>
        </tr>
    </table>
    
    <div class="content">
        <p>这是div中的内容。</p>
    </div>
</body>
</html>"""

        test_file = self.create_test_file(test_content, "test.html")

        # 测试解析
        parser = HtmlParser()
        result = parser(test_file)

        # 验证结果
        self.assertIsNotNone(result)
        self.assertIsInstance(result, list)
        self.assertGreater(len(result), 0)

        # 检查内容
        all_text = " ".join(result)
        self.assertIn("主标题", all_text)
        self.assertIn("这是div中的内容。", all_text)

        print(f"✅ HTML解析成功: {len(result)} 行")

    def test_excel_parser(self):
        """测试Excel解析器"""
        # 创建简单的CSV文件作为Excel测试
        test_content = """姓名,年龄,城市,职业
张三,25,北京,工程师
李四,30,上海,设计师
王五,28,广州,产品经理
赵六,35,深圳,销售经理"""

        test_file = self.create_test_file(test_content, "test.csv")

        # 读取文件为二进制
        with open(test_file, "rb") as f:
            binary_content = f.read()

        # 测试解析
        parser = ExcelParser()
        result = parser(binary_content)

        # 验证结果
        self.assertIsNotNone(result)
        self.assertIsInstance(result, list)
        self.assertGreater(len(result), 0)

        # 检查内容
        all_text = " ".join(result)
        self.assertIn("张三", all_text)
        self.assertIn("工程师", all_text)

        print(f"✅ Excel/CSV解析成功: {len(result)} 行")

    def test_excel_parser_html_output(self):
        """测试Excel解析器的HTML输出"""
        test_content = """姓名,年龄,城市,职业
张三,25,北京,工程师
李四,30,上海,设计师"""

        test_file = self.create_test_file(test_content, "test_html.csv")

        # 读取文件为二进制
        with open(test_file, "rb") as f:
            binary_content = f.read()

        # 测试HTML输出
        parser = ExcelParser()
        result = parser.html(binary_content)

        # 验证结果
        self.assertIsNotNone(result)
        self.assertIsInstance(result, list)
        self.assertGreater(len(result), 0)

        # 检查HTML内容
        html_content = result[0]
        self.assertIn("<table>", html_content)
        self.assertIn("<th>", html_content)
        self.assertIn("张三", html_content)

        print(f"✅ Excel HTML输出成功: {len(result)} 个表格")

    def test_ppt_parser(self):
        """测试PPT解析器"""
        test_file = self.create_test_pptx_file("test.pptx")

        parser = PptParser()
        result = parser(test_file, from_page=0, to_page=10)

        self.assertIsNotNone(result)
        self.assertIsInstance(result, list)
        self.assertGreaterEqual(len(result), 2)

        all_text = " ".join(result)
        self.assertIn("测试PPT标题", all_text)
        self.assertIn("第一页内容", all_text)
        self.assertIn("第二页内容", all_text)

        print(f"✅ PPT解析成功: {len(result)} 页")

    def test_docx_parser(self):
        """测试Word解析器"""
        test_file = self.create_test_docx_file("test.docx")

        parser = DocxParser()
        secs, tbls = parser(test_file)

        self.assertIsInstance(secs, list)
        self.assertIsInstance(tbls, list)
        self.assertGreater(len(secs), 0)

        all_text = " ".join(text for text, _style in secs if text)
        self.assertIn("测试Word标题", all_text)
        self.assertIn("这是Word正文内容。", all_text)

        all_tables = " ".join(line for table in tbls for line in table)
        self.assertIn("列1: 数据1", all_tables)
        self.assertIn("列2: 数据2", all_tables)

        print(f"✅ DOCX解析成功: {len(secs)} 个段落, {len(tbls)} 个表格")

    def test_pdf_parser(self):
        """测试PDF解析器"""
        test_file = self.create_test_pdf_file("test.pdf")

        parser = PdfParser()
        try:
            with warnings.catch_warnings():
                warnings.simplefilter("ignore")
                text_content, tbls = parser(test_file, need_image=False)
        except Exception as exc:
            msg = str(exc).lower()
            asset_markers = ("model", "onnx", "xgb", "nltk", "wordnet", "huqie", "no such file", "not found")
            if any(marker in msg for marker in asset_markers):
                self.skipTest(f"PdfParser runtime assets unavailable: {exc}")
            raise

        self.assertIsInstance(text_content, str)
        self.assertIsInstance(tbls, list)
        self.assertEqual(parser.total_page, 1)

        print(f"✅ PDF解析成功: {len(text_content)} 个文本字符, {len(tbls)} 个表格/图片块")

    @unittest.skip("需要真实的图片文件")
    def test_figure_parser(self):
        """测试图片解析器"""
        # 这个测试需要真实的图片文件
        pass

    def test_parser_error_handling(self):
        """测试解析器错误处理"""
        # 测试不存在的文件
        non_existent_file = self.test_data_dir / "non_existent.txt"

        parser = TxtParser()
        with self.assertRaises(FileNotFoundError):
            parser(os.fspath(non_existent_file))

        print("✅ 错误处理测试通过")

    def test_encoding_handling(self):
        """测试编码处理"""
        # 创建GBK编码的文件
        test_content = "这是GBK编码的中文内容"
        test_file = self.create_test_file(test_content, "test_gbk.txt", encoding="gbk")

        parser = TxtParser()
        with open(test_file, "rb") as f:
            binary_content = f.read()

        result = parser(test_file, binary_content)

        self.assertIsNotNone(result)
        self.assertIsInstance(result, list)

        # 检查内容
        all_text = " ".join([chunk[0] for chunk in result])
        self.assertIn("GBK编码", all_text)

        print("✅ 编码处理测试通过")


class TestParserIntegration(unittest.TestCase):
    """测试解析器集成功能"""

    def setUp(self):
        """测试前的准备工作"""
        self.test_data_dir = Path(__file__).parent / "test_data"
        self.test_data_dir.mkdir(exist_ok=True)

    def test_multiple_parsers(self):
        """测试多个解析器同时工作"""
        # 创建不同类型的测试文件
        files = {
            "test1.txt": "这是文本文件内容",
            "test2.md": "# 标题\n这是Markdown内容",
            "test3.json": '{"name": "测试", "value": 123}',
            "test4.html": "<html><body><h1>标题</h1><p>内容</p></body></html>",
            "test5.csv": "姓名,年龄\n张三,25\n李四,30",
        }

        parsers = {"txt": TxtParser(), "md": MarkdownParser(), "json": JsonParser(), "html": HtmlParser(), "excel": ExcelParser()}

        results = {}

        for filename, content in files.items():
            file_path = self.test_data_dir / filename
            with open(file_path, "w", encoding="utf-8") as f:
                f.write(content)

            # 根据文件扩展名选择解析器
            ext = filename.split(".")[-1]
            if ext == "txt":
                parser = parsers["txt"]
                result = parser(os.fspath(file_path))
            elif ext == "md":
                parser = parsers["md"]
                with open(file_path, "r", encoding="utf-8") as f:
                    content = f.read()
                text, tables = parser.extract_tables_and_remainder(content)
                result = [text]  # 简化处理
            elif ext == "json":
                parser = parsers["json"]
                with open(file_path, "rb") as f:
                    binary_content = f.read()
                result = parser(binary_content)
            elif ext == "html":
                parser = parsers["html"]
                result = parser(os.fspath(file_path))
            elif ext == "csv":
                parser = parsers["excel"]
                with open(file_path, "rb") as f:
                    binary_content = f.read()
                result = parser(binary_content)
            else:
                continue

            results[filename] = result

        # 验证所有解析都成功
        self.assertEqual(len(results), 5)
        for filename, result in results.items():
            self.assertIsNotNone(result)
            self.assertIsInstance(result, list)
            self.assertGreater(len(result), 0)

        print("✅ 多解析器集成测试通过")


if __name__ == "__main__":
    # 运行所有测试
    unittest.main(verbosity=2)
